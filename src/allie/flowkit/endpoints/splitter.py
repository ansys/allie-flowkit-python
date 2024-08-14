# Copyright (C) 2024 ANSYS, Inc. and/or its affiliates.
# SPDX-License-Identifier: MIT
#
#
# Permission is hereby granted, free of charge, to any person obtaining a copy
# of this software and associated documentation files (the "Software"), to deal
# in the Software without restriction, including without limitation the rights
# to use, copy, modify, merge, publish, distribute, sublicense, and/or sell
# copies of the Software, and to permit persons to whom the Software is
# furnished to do so, subject to the following conditions:
#
# The above copyright notice and this permission notice shall be included in all
# copies or substantial portions of the Software.
#
# THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY KIND, EXPRESS OR
# IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF MERCHANTABILITY,
# FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT. IN NO EVENT SHALL THE
# AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER
# LIABILITY, WHETHER IN AN ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING FROM,
# OUT OF OR IN CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER DEALINGS IN THE
# SOFTWARE.

"""Module for splitting text into chunks."""

import base64
import io

from allie.flowkit.config._config import CONFIG
from allie.flowkit.models.splitter import SplitterRequest, SplitterResponse
from fastapi import APIRouter, Header, HTTPException
from langchain.text_splitter import PythonCodeTextSplitter, RecursiveCharacterTextSplitter
from pdfminer.high_level import extract_text
from pptx import Presentation

TOKEN_TO_CHARACTER_MULTIPLIER = 4

router = APIRouter()


@router.post("/ppt", response_model=SplitterResponse)
async def split_ppt(request: SplitterRequest, api_key: str = Header(...)) -> SplitterResponse:
    """Endpoint for splitting text in a PowerPoint document into chunks.

    Parameters
    ----------
    request : SplitterRequest
        An object containing 'document_content' in Base64,
        'chunk_size', and 'chunk_overlap'
    api_key : str
        The API key for authentication.

    """
    validate_request(request, api_key)
    return process_ppt(request)


@router.post("/py", response_model=SplitterResponse)
async def split_py(request: SplitterRequest, api_key: str = Header(...)) -> SplitterResponse:
    """Endpoint for splitting Python code into chunks.

    Parameters
    ----------
    request : SplitterRequest
        An object containing 'document_content' in Base64,
        'chunk_size', and 'chunk_overlap'
    api_key : str
        The API key for authentication.

    Returns
    -------
    SplitterResponse
        An object containing a list of text chunks.

    """
    validate_request(request, api_key)
    return process_python_code(request)


@router.post("/pdf", response_model=SplitterResponse)
async def split_pdf(request: SplitterRequest, api_key: str = Header(...)) -> SplitterResponse:
    """Endpoint for splitting text in a PDF document into chunks.

    Parameters
    ----------
    request : SplitterRequest
        An object containing 'document_content' in Base64,
        'chunk_size', and 'chunk_overlap'.
    api_key : str
        The API key for authentication.

    Returns
    -------
    SplitterResponse
        An object containing a list of text chunks.

    """
    validate_request(request, api_key)
    return process_pdf(request)


def process_ppt(request: SplitterRequest) -> SplitterResponse:
    """Process a PowerPoint document to split text into chunks.

    Parameters
    ----------
    request : SplitterRequest
        An object containing 'document_content' in Base64,
        'chunk_size', and 'chunk_overlap'

    Returns
    -------
    SplitterResponse
        An object containing a list of text chunks.

    """
    try:
        document_content = base64.b64decode(request.document_content)
    except base64.binascii.Error:
        raise HTTPException(status_code=400, detail="Invalid Base64 encoding")

    try:
        ppt_document = Presentation(io.BytesIO(document_content))
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error processing PowerPoint file: {str(e)}")

    ppt_text = ""
    for slide in ppt_document.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        ppt_text += run.text + " "

    if not ppt_text:
        raise HTTPException(status_code=400, detail="No text found in PowerPoint document")

    chunk_size_langchain = request.chunk_size * TOKEN_TO_CHARACTER_MULTIPLIER
    chunk_overlap_langchain = request.chunk_overlap * TOKEN_TO_CHARACTER_MULTIPLIER

    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size_langchain, chunk_overlap=chunk_overlap_langchain)
    chunks = splitter.split_text(ppt_text)
    response = SplitterResponse(chunks=chunks)

    return response


def process_python_code(request: SplitterRequest) -> SplitterResponse:
    """Process Python code to split text into chunks.

    Parameters
    ----------
    request : SplitterRequest
        An object containing 'document_content' in Base64,
        'chunk_size', and 'chunk_overlap'

    Returns
    -------
    SplitterResponse
        An object containing a list of text chunks.

    """
    try:
        document_content = base64.b64decode(request.document_content)
    except base64.binascii.Error:
        raise HTTPException(status_code=400, detail="Invalid Base64 encoding")

    try:
        document_content_str = document_content.decode("utf-8")
    except UnicodeDecodeError:
        raise HTTPException(status_code=400, detail="Error decoding Python code")

    chunk_size_langchain = request.chunk_size * TOKEN_TO_CHARACTER_MULTIPLIER
    chunk_overlap_langchain = request.chunk_overlap * TOKEN_TO_CHARACTER_MULTIPLIER

    splitter = PythonCodeTextSplitter(chunk_size=chunk_size_langchain, chunk_overlap=chunk_overlap_langchain)
    chunks = splitter.split_text(document_content_str)
    response = SplitterResponse(chunks=chunks)

    return response


def process_pdf(request: SplitterRequest) -> SplitterResponse:
    """Process a PDF document to split text into chunks.

    Parameters
    ----------
    request : SplitterRequest
        An object containing 'document_content' in Base64,
        'chunk_size', and 'chunk_overlap'

    Returns
    -------
    SplitterResponse
        An object containing a list of text chunks.

    """
    try:
        document_content = base64.b64decode(request.document_content)
    except base64.binascii.Error:
        raise HTTPException(status_code=400, detail="Invalid Base64 encoding")

    try:
        pdf_text = extract_text(io.BytesIO(document_content))
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error processing PDF file: {str(e)}")
    
    if not pdf_text:
        raise HTTPException(status_code=400, detail="No text found in PDF document")

    chunk_size_langchain = request.chunk_size * TOKEN_TO_CHARACTER_MULTIPLIER
    chunk_overlap_langchain = request.chunk_overlap * TOKEN_TO_CHARACTER_MULTIPLIER

    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size_langchain, chunk_overlap=chunk_overlap_langchain)
    chunks = splitter.split_text(pdf_text)
    response = SplitterResponse(chunks=chunks)

    return response


def validate_request(request: SplitterRequest, api_key: str):
    """Validate the splitter request and API key.

    Parameters
    ----------
    request : SplitterRequest
        An object containing 'document_content' in Base64,
        'chunk_size', and 'chunk_overlap'
    api_key : str
        The API key for authentication.

    Raises
    ------
    HTTPException
        If the API key is invalid or if any of the request parameters are invalid.

    """
    # Check if the provided API key matches the expected API key
    if api_key != CONFIG.flowkit_python_api_key:
        raise HTTPException(status_code=401, detail="Invalid API key")

    # Check if document content is provided
    if not request.document_content:
        raise HTTPException(status_code=400, detail="No document content provided")

    # Check if chunk size is provided
    if not request.chunk_size:
        raise HTTPException(status_code=400, detail="No chunk size provided")

    # Check if chunk size is greater than 0
    if request.chunk_size <= 0:
        raise HTTPException(status_code=400, detail="Chunk size must be greater than 0")

    # Check if chunk overlap is greater than or equal to 0
    if request.chunk_overlap < 0:
        raise HTTPException(status_code=400, detail="Chunk overlap must be greater than or equal to 0")
