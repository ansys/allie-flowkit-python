import base64
from fastapi import APIRouter
from langchain.text_splitter import RecursiveCharacterTextSplitter, PythonCodeTextSplitter
from pptx import Presentation
import io
import pymupdf
from app.models.splitter import SplitterRequest, SplitterResponse

router = APIRouter()

@router.post('/ppt', response_model=SplitterResponse)
async def split_ppt(request: SplitterRequest) -> SplitterResponse:
    """Endpoint for splitting text in a PowerPoint document into chunks.

    Parameters
    ----------
    request : SplitterRequest
        An object containing 'document_content' in Base64, 'chunk_size', and 'chunk_overlap'.
    """
    return process_ppt(request)

@router.post('/py', response_model=SplitterResponse)
async def split_py(request: SplitterRequest) -> SplitterResponse:
    """Endpoint for splitting Python code into chunks.

    Parameters
    ----------
    request : SplitterRequest
        An object containing 'document_content' in Base64, 'chunk_size', and 'chunk_overlap'.
        
    Returns
    -------
    SplitterResponse
        An object containing a list of text chunks.
    """
    return process_python_code(request)

@router.post('/pdf', response_model=SplitterResponse)
async def split_pdf(request: SplitterRequest) -> SplitterResponse:
    """Endpoint for splitting text in a PDF document into chunks.

    Parameters
    ----------
    request : SplitterRequest
        An object containing 'document_content' in Base64, 'chunk_size', and 'chunk_overlap'.
    
    Returns
    -------
    SplitterResponse
        An object containing a list of text chunks.
    """
    return process_pdf(request)


def process_ppt(request: SplitterRequest) -> SplitterResponse:
    """Process a PowerPoint document to split text into chunks.

    Parameters
    ----------
    request : SplitterRequest
        An object containing 'document_content' in Base64, 'chunk_size', and 'chunk_overlap'.
    
    Returns
    -------
    SplitterResponse
        An object containing a list of text chunks.
    """
    document_content = base64.b64decode(request.document_content)
    ppt_document = Presentation(io.BytesIO(document_content))
    ppt_text = ""

    for slide in ppt_document.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        ppt_text += run.text + " "

    splitter = RecursiveCharacterTextSplitter(chunk_size=request.chunk_size, chunk_overlap=request.chunk_overlap)
    chunks = splitter.split_text(ppt_text)
    response = SplitterResponse(chunks=chunks)
    
    return response

def process_python_code(request: SplitterRequest) -> SplitterResponse:
    """Process Python code to split text into chunks.

    Parameters
    ----------
    request : SplitterRequest
        An object containing 'document_content' in Base64, 'chunk_size', and 'chunk_overlap'.
    
    Returns
    -------
    SplitterResponse
        An object containing a list of text chunks.
    """
    document_content = base64.b64decode(request.document_content)
    document_content_str = document_content.decode('utf-8')

    splitter = PythonCodeTextSplitter(chunk_size=request.chunk_size, chunk_overlap=request.chunk_overlap)
    chunks = splitter.split_text(document_content_str)
    response = SplitterResponse(chunks=chunks)
    
    return response

def process_pdf(request: SplitterRequest) -> SplitterResponse:
    """Process a PDF document to split text into chunks.

    Parameters
    ----------
    request : SplitterRequest
        An object containing 'document_content' in Base64, 'chunk_size', and 'chunk_overlap'.
    
    Returns
    -------
    SplitterResponse
        An object containing a list of text chunks.
    """
    document_content = base64.b64decode(request.document_content)
    pdf_document = pymupdf.open(stream=io.BytesIO(document_content), filetype="pdf")
    pdf_text = ""

    for page in pdf_document:
        pdf_text += page.get_text()

    splitter = RecursiveCharacterTextSplitter(chunk_size=request.chunk_size, chunk_overlap=request.chunk_overlap)
    chunks = splitter.split_text(pdf_text)
    response = SplitterResponse(chunks=chunks)
    
    return response